import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from tabulate import tabulate
from tqdm import tqdm
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


# 遍历目录树，获取所有子目录下的图片文件路径
def get_image_files(directory):
    image_files = []
    for root, dirs, files in os.walk(directory):
        for file in files:
            if file.endswith(".png") or file.endswith(".jpg"):
                image_files.append(os.path.join(root, file))
    return image_files




def create_ppt(directory, filenames):
    prs = Presentation('root\\result\\visualization\\data summary.pptx')
    width = Inches(10)
    height = Inches(7.5)
        
    # 在第二页添加目录
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    title_shape = slide.shapes.title
    title_shape.text = "目录"
    body_shape = slide.shapes.placeholders[1]
# 在每个子目录的第一张幻灯片上添加链接和页码
    for subdir in os.listdir(directory):
        subdir_path = os.path.join(directory, subdir)
        if os.path.isdir(subdir_path):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = subdir
            slide.shapes.title.anchor = "start"
            slide.shapes.title.id = "{}".format(subdir)
            image_files = get_image_files(subdir_path)
            for i, image_file in enumerate(tqdm(image_files, desc=subdir)):
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                img_path = image_file
                with Image.open(img_path) as img:
                    pic = slide.shapes.add_picture(img_path, Inches(0.5), Inches(0.8), width=0.8*width, height=0.55*height)
                slide.shapes.title.text = os.path.splitext(os.path.basename(image_file))[0]
                textbox = slide.shapes.placeholders[0]
                textbox.text = os.path.splitext(os.path.basename(image_file))[0]
                textbox.text_frame.paragraphs[0].font.name = '宋体'
                textbox.text_frame.paragraphs[0].font.size = Pt(14)
                textbox.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
                
                
                # 添加页码
                text_box = slide.shapes.add_textbox(Inches(9), Inches(6.5), Inches(1), Inches(0.5))
                text_box.text_frame.text = str(i+1)
                text_box.text_frame.paragraphs[0].font.size = Pt(14)
                text_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
                text_box.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    prs.save('data summary.pptx')

create_ppt("root\\result\\visualization", "filenames")

