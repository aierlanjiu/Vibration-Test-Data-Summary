from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches
import os

# Open the PowerPoint file
ppt = Presentation('root\\result\\visualization\\data summary.pptx')

# Find slides containing target line
target_slides = []
for slide in ppt.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                if 'run up' in paragraph.text and 'target line' in paragraph.text:
                    target_slides.append((slide, 'run up'))
                    break
                elif 'run down' in paragraph.text and 'target line' in paragraph.text:
                    target_slides.append((slide, 'run down'))
                    break

# Delete images on target slides
for slide, position in target_slides:
    for shape in slide.shapes:
        shapes = slide.shapes
        for i in range(len(shapes)-1, -1, -1):
            shape = shapes[i]
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                shapes._spTree.remove(shape.element)

# Insert images from a folder
images = []
for filename in os.listdir('08 dataoverview\\dataoverview\\04 targetline'):
    if filename.endswith('.png'):
        images.append(filename)

# Insert images into target slides
height1, width1 = Inches(1.57), Inches(2.95) # size of images for run up
height2, width2 = Inches(2), Inches(3) # size of images for run down
count1, count2 = 0, 36 # count of images inserted into run up and run down slides
for slide, position in target_slides:
    if position == 'run up':
        for i in range(2):
            for j in range(3):
                if count1 < len(images):
                    left = j * width2 + Inches(0)
                    top = i * height2 + Inches(1)
                    slide.shapes.add_picture('08 dataoverview/dataoverview/04 targetline/' + images[count1], left, top, width=width2, height=height2)
                    count1 += 1
    elif position == 'run down':
        for i in range(2):
            for j in range(3):
                if count2 < len(images):
                    left = j * width2 + Inches(0)
                    top = i * height2 + Inches(1)
                    slide.shapes.add_picture('08 dataoverview/dataoverview/04 targetline/' + images[count2], left, top, width=width2, height=height2)
                    count2 += 1



# Save the PowerPoint file
ppt.save('root\\result\\visualization\\data summary_V0.pptx')